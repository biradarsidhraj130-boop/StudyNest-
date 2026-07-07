"""
Universal Document Ingestion & OCR Fallback
Modular text extraction system for multiple document types with OCR fallback
"""

import os
import mimetypes
from pathlib import Path
from typing import Any, Dict, Optional
from dataclasses import dataclass
import logging

# Configure logging
logger = logging.getLogger(__name__)

@dataclass
class ExtractResult:
    """Result of text extraction from a document"""
    text: str
    metadata: Dict[str, Any]
    success: bool
    error_message: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for storage"""
        return {
            "text": self.text,
            "metadata": self.metadata,
            "success": self.success,
            "error_message": self.error_message
        }

class DocumentIngester:
    """Modular document text extraction system"""
    
    def __init__(self, enable_ocr: bool = True):
        self.enable_ocr = enable_ocr and self._check_ocr_available()
        self.extractors = {
            # Existing formats
            ".pdf": self._extract_pdf_text,
            ".docx": self._extract_docx_text,
            ".txt": self._extract_txt_text,
            
            # New text-based formats
            ".pptx": self._extract_pptx_text,
            ".xlsx": self._extract_xlsx_text,
            ".html": self._extract_html_text,
            ".htm": self._extract_html_text,
            
            # OCR-based formats
            ".png": self._extract_image_ocr,
            ".jpg": self._extract_image_ocr,
            ".jpeg": self._extract_image_ocr,
        }
    
    def _check_ocr_available(self) -> bool:
        """Check if OCR dependencies are available"""
        try:
            import pytesseract
            from PIL import Image
            import pdf2image
            return True
        except ImportError as e:
            logger.warning(f"OCR dependencies not available: {e}")
            return False
    
    def extract_text_and_metadata(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """
        Extract text and metadata from a document using appropriate extractor
        
        Args:
            file_path: Path to the uploaded file
            filename: Original filename
            mime_type: Optional MIME type for format detection
            
        Returns:
            ExtractResult with extracted text, metadata, and success status
        """
        try:
            # Determine file extension
            extension = Path(filename).suffix.lower()
            
            # Get MIME type if not provided
            if not mime_type:
                mime_type, _ = mimetypes.guess_type(filename)
            
            # Select appropriate extractor
            extractor = self.extractors.get(extension)
            
            if not extractor:
                return ExtractResult(
                    text="",
                    metadata={
                        "file_type": extension,
                        "mime_type": mime_type,
                        "extractor": "none",
                        "unreadable": True
                    },
                    success=False,
                    error_message=f"Unsupported file type: {extension}. Supported types: {', '.join(self.extractors.keys())}"
                )
            
            # Extract text using appropriate method
            result = extractor(file_path, filename, mime_type)
            
            # Add common metadata
            result.metadata.update({
                "file_type": extension,
                "mime_type": mime_type,
                "extractor": extractor.__name__.replace("_extract_", "").replace("_text", ""),
                "file_size": file_path.stat().st_size
            })
            
            return result
            
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            return ExtractResult(
                text="",
                metadata={
                    "file_type": Path(filename).suffix.lower(),
                    "mime_type": mime_type,
                    "extractor": "error",
                    "unreadable": True
                },
                success=False,
                error_message=f"Extraction failed: {str(e)}"
            )
    
    def _extract_pdf_text(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from PDF with OCR fallback for scanned documents"""
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            text = ""
            page_count = doc.page_count
            
            # Try to extract text directly first
            for page in doc:
                page_text = page.get_text()
                text += page_text + "\n"
            
            doc.close()
            
            # Check if extracted text is meaningful (not just scanned)
            text_stripped = text.strip()
            if len(text_stripped) < 50:  # Very little text, likely scanned
                if self.enable_ocr:
                    logger.info(f"PDF appears to be scanned, attempting OCR: {filename}")
                    return self._extract_scanned_pdf_ocr(file_path, filename, mime_type)
                else:
                    return ExtractResult(
                        text="",
                        metadata={
                            "page_count": page_count,
                            "scanned": True,
                            "ocr_used": False,
                            "unreadable": True
                        },
                        success=False,
                        error_message="PDF appears to be scanned but OCR is disabled or unavailable"
                    )
            
            return ExtractResult(
                text=text_stripped,
                metadata={
                    "page_count": page_count,
                    "scanned": False,
                    "ocr_used": False,
                    "text_length": len(text_stripped)
                },
                success=True
            )
            
        except ImportError:
            return ExtractResult(
                text="",
                metadata={"extractor": "pdf", "missing_dep": "PyMuPDF"},
                success=False,
                error_message="PDF extraction requires PyMuPDF (fitz) package"
            )
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={"extractor": "pdf", "error": str(e)},
                success=False,
                error_message=f"PDF extraction failed: {str(e)}"
            )
    
    def _extract_docx_text(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from DOCX documents"""
        try:
            import docx
            
            doc = docx.Document(file_path)
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
            
            return ExtractResult(
                text=text.strip(),
                metadata={
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(doc.tables),
                    "text_length": len(text)
                },
                success=True
            )
            
        except ImportError:
            return ExtractResult(
                text="",
                metadata={"extractor": "docx", "missing_dep": "python-docx"},
                success=False,
                error_message="DOCX extraction requires python-docx package"
            )
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={"extractor": "docx", "error": str(e)},
                success=False,
                error_message=f"DOCX extraction failed: {str(e)}"
            )
    
    def _extract_txt_text(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            text = ""
            encoding_used = None
            
            for encoding in encodings:
                try:
                    text = file_path.read_text(encoding=encoding)
                    encoding_used = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if not encoding_used:
                # Last resort with error ignoring
                text = file_path.read_bytes().decode('utf-8', errors='ignore')
                encoding_used = 'utf-8-ignore'
            
            return ExtractResult(
                text=text.strip(),
                metadata={
                    "encoding": encoding_used,
                    "text_length": len(text)
                },
                success=True
            )
            
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={"extractor": "txt", "error": str(e)},
                success=False,
                error_message=f"Text extraction failed: {str(e)}"
            )
    
    def _extract_pptx_text(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from PowerPoint presentations"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    # Extract text from tables
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                text += cell.text + " "
            
            return ExtractResult(
                text=text.strip(),
                metadata={
                    "slides": slide_count,
                    "text_length": len(text)
                },
                success=True
            )
            
        except ImportError:
            return ExtractResult(
                text="",
                metadata={"extractor": "pptx", "missing_dep": "python-pptx"},
                success=False,
                error_message="PPTX extraction requires python-pptx package"
            )
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={"extractor": "pptx", "error": str(e)},
                success=False,
                error_message=f"PPTX extraction failed: {str(e)}"
            )
    
    def _extract_xlsx_text(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from Excel spreadsheets"""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            text = ""
            sheet_count = 0
            cell_count = 0
            
            for sheet_name in workbook.sheetnames:
                sheet_count += 1
                sheet = workbook[sheet_name]
                
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            text += str(cell) + " "
                            cell_count += 1
            
            workbook.close()
            
            return ExtractResult(
                text=text.strip(),
                metadata={
                    "sheets": sheet_count,
                    "cells_with_data": cell_count,
                    "text_length": len(text)
                },
                success=True
            )
            
        except ImportError:
            return ExtractResult(
                text="",
                metadata={"extractor": "xlsx", "missing_dep": "openpyxl"},
                success=False,
                error_message="XLSX extraction requires openpyxl package"
            )
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={"extractor": "xlsx", "error": str(e)},
                success=False,
                error_message=f"XLSX extraction failed: {str(e)}"
            )
    
    def _extract_html_text(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from HTML files"""
        try:
            from bs4 import BeautifulSoup
            
            html_content = file_path.read_text(encoding='utf-8', errors='ignore')
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return ExtractResult(
                text=text,
                metadata={
                    "title": soup.title.string if soup.title else None,
                    "text_length": len(text)
                },
                success=True
            )
            
        except ImportError:
            # Fallback to basic HTML parsing without BeautifulSoup
            try:
                import re
                html_content = file_path.read_text(encoding='utf-8', errors='ignore')
                
                # Remove HTML tags
                text = re.sub(r'<[^>]+>', ' ', html_content)
                text = re.sub(r'\s+', ' ', text).strip()
                
                return ExtractResult(
                    text=text,
                    metadata={
                        "extractor": "html_regex",
                        "text_length": len(text)
                    },
                    success=True
                )
            except Exception as e:
                return ExtractResult(
                    text="",
                    metadata={"extractor": "html", "missing_dep": "beautifulsoup4", "fallback_error": str(e)},
                    success=False,
                    error_message="HTML extraction recommends beautifulsoup4 package"
                )
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={"extractor": "html", "error": str(e)},
                success=False,
                error_message=f"HTML extraction failed: {str(e)}"
            )
    
    def _extract_image_ocr(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from images using OCR"""
        if not self.enable_ocr:
            return ExtractResult(
                text="",
                metadata={
                    "ocr_used": False,
                    "ocr_disabled": True,
                    "unreadable": True
                },
                success=False,
                error_message="OCR is disabled or unavailable. Install pytesseract, Pillow, and system Tesseract."
            )
        
        try:
            from PIL import Image
            import pytesseract
            
            # Open image
            image = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            return ExtractResult(
                text=text.strip(),
                metadata={
                    "ocr_used": True,
                    "image_size": image.size,
                    "image_mode": image.mode,
                    "text_length": len(text.strip())
                },
                success=bool(text.strip())
            )
            
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={
                    "ocr_used": False,
                    "ocr_error": str(e)
                },
                success=False,
                error_message=f"OCR extraction failed: {str(e)}"
            )
    
    def _extract_scanned_pdf_ocr(self, file_path: Path, filename: str, mime_type: Optional[str] = None) -> ExtractResult:
        """Extract text from scanned PDFs using OCR"""
        if not self.enable_ocr:
            return ExtractResult(
                text="",
                metadata={
                    "ocr_used": False,
                    "ocr_disabled": True,
                    "unreadable": True
                },
                success=False,
                error_message="OCR is disabled or unavailable for scanned PDF processing."
            )
        
        try:
            import pdf2image
            from PIL import Image
            import pytesseract
            
            # Convert PDF to images
            images = pdf2image.convert_from_path(file_path)
            text = ""
            
            for i, image in enumerate(images):
                page_text = pytesseract.image_to_string(image)
                text += f"Page {i+1}:\n{page_text}\n"
            
            return ExtractResult(
                text=text.strip(),
                metadata={
                    "ocr_used": True,
                    "pages_processed": len(images),
                    "text_length": len(text.strip())
                },
                success=bool(text.strip())
            )
            
        except Exception as e:
            return ExtractResult(
                text="",
                metadata={
                    "ocr_used": False,
                    "ocr_error": str(e)
                },
                success=False,
                error_message=f"Scanned PDF OCR failed: {str(e)}"
            )

# Global instance for easy import
def create_document_ingester() -> DocumentIngester:
    """Create a document ingester instance with OCR settings from environment"""
    enable_ocr = os.getenv("ENABLE_OCR", "true").lower() == "true"
    return DocumentIngester(enable_ocr=enable_ocr)
